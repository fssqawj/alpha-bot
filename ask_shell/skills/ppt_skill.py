"""PPT Generation Skill - Create PowerPoint presentations"""

import os
import re
import json
from loguru import logger
from datetime import datetime
from typing import List, Optional, Dict, Any, Callable
from .base_skill import BaseSkill, SkillExecutionResponse, SkillCapability
from ..llm.base import BaseLLMClient
from ..llm.openai_client import OpenAIClient

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    Presentation = None


class PPTSkill(BaseSkill):
    """
    Skill for generating PowerPoint presentations
    
    This skill can:
    1. Create PPT from text outlines
    2. Generate slides with specific themes
    3. Add charts and images to presentations
    """
    
    def __init__(self):
        super().__init__()
        self.initialized = Presentation is not None
        self.output_dir = "./output"
        if not os.path.exists(self.output_dir):
            os.makedirs(self.output_dir, exist_ok=True)
        
        # Initialize LLM client for content generation
        try:
            self.llm: BaseLLMClient = OpenAIClient()
            self.system_prompt = """你是一个专业的PPT内容策划师。用户会给你一个主题和任务要求，以及可能的历史交互信息。请为PowerPoint演示文稿生成合适的大纲和每页的详细内容。

你的回复必须是一个JSON对象，格式如下：
{
    "title": "演示文稿的整体标题",
    "outline": [
        {
            "title": "幻灯片标题",
            "content": "幻灯片详细内容，包括要点、说明等"
        },
        ...
    ]
}

要求：
1. 幻灯片数量通常为3-8张，根据内容复杂程度调整
2. 内容要有层次感，从概述到细节逐步展开
3. 考虑历史信息，使内容连贯一致
4. 确保内容专业、有深度且易于理解
5. 适应用户的具体需求和任务背景
6. **美化要求**：内容应适合美观的PPT展示，包含清晰的标题、要点分明的内容，适合视觉呈现，考虑使用列表、要点、短句等形式便于PPT美化排版"""
            self.llm.set_system_prompt(self.system_prompt)
            self.llm.set_direct_mode(True)  # Use direct mode for content generation
            self.llm_available = True
        except Exception as e:
            # If LLM initialization fails (e.g., missing API key), still allow basic functionality
            print(f"Warning: LLM initialization failed: {str(e)}. PPT skill will use basic functionality.")
            self.llm_available = False
    
    def get_capabilities(self) -> List[SkillCapability]:
        """PPT skill provides file generation capability"""
        return [SkillCapability.FILE_GENERATION]
    
    def execute(
        self,
        task: str,
        context: Optional[Dict[str, Any]] = None,
        **kwargs
    ) -> SkillExecutionResponse:
        """
        Execute PPT generation
            
        Args:
            task: Task description (e.g., "create a PPT about AI")
            context: Execution context
            **kwargs: Additional parameters including selection_reasoning
            
        Returns:
            SkillExecutionResponse with generated file path
        """
        # Get the reasoning for why this skill was selected
        selection_reasoning = kwargs.get('selection_reasoning', '')
        
        if not self.initialized:
            return SkillExecutionResponse(
                thinking="Checking if required library (python-pptx) is installed",
                direct_response="Error: python-pptx library is not installed. Please install it using: pip install python-pptx",
                generated_files=[],
                file_metadata={"status": "missing_dependency", "required_library": "python-pptx"}
            )
        
        try:
            # Get execution context from context
            last_result = context.get('last_result')
            history = context.get('history', [])
            
            # Generate presentation outline using LLM based on task and context
            title, outline = self._generate_outline_with_llm(task, history, last_result)
            
            # Generate the presentation
            filename = self._generate_presentation(title, outline)
            
            return SkillExecutionResponse(
                thinking=f"Creating PowerPoint presentation '{title}' with {len(outline)} slides based on the given task and context",
                direct_response=f"Successfully created PowerPoint presentation '{filename}' with {len(outline)} slides.",
                generated_files=[filename],
                file_metadata={
                    "status": "success",
                    "slides_count": len(outline),
                    "title": title,
                    "filename": filename
                }
            )
        except Exception as e:
            return SkillExecutionResponse(
                thinking=f"Failed to create PowerPoint presentation: {str(e)}",
                direct_response=f"Error creating PowerPoint presentation: {str(e)}",
                generated_files=[],
                file_metadata={"status": "error", "error": str(e)}
            )
    
    def _generate_outline_with_llm(self, task: str, history: List[Any], last_result: Optional[Any]) -> tuple[str, List[Dict[str, Any]]]:
        """Generate presentation outline using LLM based on task and context"""
        
        # If LLM is not available, fall back to basic parsing
        if not self.llm_available:
            return self._parse_task_basic(task, history, last_result)
        
        context_str = f"{task}\n\n"
        
        message = f"{context_str}\n\n请根据以上信息为PowerPoint演示文稿生成合适的大纲和每页的详细内容。" 
        
        try:
            # Call LLM to generate outline with direct parsing using PPTSkillResponse dataclass
            logger.info(f"PPT Skill LLM Prompt: {message}")
            
            from ..models.types import PPTSkillResponse
            # Generate and directly parse into PPTSkillResponse
            llm_response = self.llm.generate(message, last_result, history=history, response_class=PPTSkillResponse)
            logger.info(f"PPT Skill LLM Response: {llm_response}")
            
            # If the response is already parsed (when response_class is provided), use it directly
            if hasattr(llm_response, 'outline'):  # It's already a PPTSkillResponse object
                parsed_response = llm_response
            else:
                # Fallback to raw JSON parsing if needed
                parsed = None
                try:
                    parsed = json.loads(llm_response.raw_json)
                except json.JSONDecodeError:
                    logger.opt(exception=True).error("Error parsing JSON from LLM response")
                    # Fallback to extracting JSON from raw response if direct parsing fails
                    response_text = llm_response.raw_json
                    start_idx = response_text.find('{')
                    end_idx = response_text.rfind('}')
                    
                    if start_idx != -1 and end_idx != -1 and start_idx < end_idx:
                        json_str = response_text[start_idx:end_idx+1]
                        try:
                            parsed = json.loads(json_str)
                        except json.JSONDecodeError:
                            logger.opt(exception=True).error("Error parsing JSON from LLM response")
                            pass
                
                if parsed:
                    # Create PPTSkillResponse manually
                    title = parsed.get("title", "Generated Presentation")
                    outline = parsed.get("outline", [])
                    parsed_response = PPTSkillResponse(title=title, outline=outline)
                else:
                    # If JSON parsing failed, use default values
                    parsed_response = PPTSkillResponse(title="Generated Presentation", outline=[])
            
            title = parsed_response.title or "Generated Presentation"
            outline = parsed_response.outline or []
            
            if title and outline:
                # Validate outline structure
                if not isinstance(outline, list) or len(outline) == 0:
                    # Fallback to a default structure
                    outline = [
                        {"title": title, "content": "Introduction slide for the presentation"},
                        {"title": "Overview", "content": "Main points of the presentation"},
                        {"title": "Details", "content": "Detailed information and analysis"},
                        {"title": "Conclusion", "content": "Summary and next steps"}
                    ]
                
                return title.strip(), outline
            
            # If JSON parsing failed, fallback to default behavior
            return "Generated Presentation", [
                {"title": "Introduction", "content": f"Presentation about: {task}"},
                {"title": "Overview", "content": "Main points of the presentation"},
                {"title": "Details", "content": "Detailed information based on context"},
                {"title": "Conclusion", "content": "Summary and next steps"}
            ]
        except Exception as e:
            logger.opt(exception=e).error("Error generating outline with LLM")
            # Fallback to default behavior if LLM call fails
            return "Generated Presentation", [
                {"title": "Introduction", "content": f"Presentation about: {task}"},
                {"title": "Overview", "content": "Main points of the presentation"},
                {"title": "Details", "content": "Detailed information based on context"},
                {"title": "Conclusion", "content": "Summary and next steps"}
            ]
    
    def _parse_task_basic(self, task: str, history: List[Any], last_result: Optional[Any]) -> tuple[str, List[Dict[str, Any]]]:
        """Basic task parsing when LLM is not available"""
        # Extract title from task
        title = task.split('about')[-1].strip() if 'about' in task.lower() else task.split('presentation')[-1].strip() if 'presentation' in task.lower() else task
        title = title or "Generated Presentation"
        
        # Create basic outline based on the task and any context from history
        outline = [
            {"title": title, "content": f"Introduction to {title}"},
            {"title": "Background", "content": "Context and importance of this topic"},
            {"title": "Key Points", "content": "Main ideas and concepts"},
            {"title": "Applications", "content": "Practical uses and examples"},
            {"title": "Future Outlook", "content": "Trends and developments"},
            {"title": "Conclusion", "content": "Summary and recommendations"}
        ]
        
        # Adjust outline based on history if available
        if history:
            context_summary = " "
            for result in history[-2:]:  # Use last 2 history items
                if hasattr(result, 'skill_response') and hasattr(result.skill_response, 'direct_response'):
                    context_summary += f" {result.skill_response.direct_response}"
                elif hasattr(result, 'command'):
                    context_summary += f" Command: {result.command}"
            
            if len(context_summary) > 2:
                outline[2]["content"] = f"Main ideas and concepts based on context:{context_summary}"
        
        # Limit to 4-6 slides for readability
        return title.strip(), outline[:6]
    
    def _generate_presentation(self, title: str, outline: List[Dict[str, Any]]) -> str:
        """Generate the actual PowerPoint presentation with enhanced styling and layouts"""
        prs = Presentation()
        
        # Apply theme and color scheme
        # Use a professional color scheme
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Pt
        
        # Add title slide with enhanced styling
        title_slide_layout = prs.slide_layouts[0]  # Title Slide layout
        title_slide = prs.slides.add_slide(title_slide_layout)
        title_placeholder = title_slide.shapes.title
        subtitle_placeholder = title_slide.placeholders[1]
        
        # Style the title slide
        if title_placeholder:
            title_placeholder.text = title
            title_frame = title_placeholder.text_frame
            if title_frame.paragraphs:
                title_para = title_frame.paragraphs[0]
                title_para.font.size = Pt(44)
                title_para.font.bold = True
                title_para.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue
                title_para.alignment = PP_ALIGN.CENTER
        
        if subtitle_placeholder:
            subtitle_placeholder.text = f"Generated on {datetime.now().strftime('%Y-%m-%d')}\nBy Ask-Shell PPT Skill"
            subtitle_frame = subtitle_placeholder.text_frame
            if subtitle_frame.paragraphs:
                subtitle_para = subtitle_frame.paragraphs[0]
                subtitle_para.font.size = Pt(18)
                subtitle_para.font.color.rgb = RGBColor(102, 102, 102)  # Gray
                subtitle_para.alignment = PP_ALIGN.CENTER
        
        # Add content slides with varied layouts and styling
        for i, slide_data in enumerate(outline):
            # Alternate between different slide layouts for visual interest
            layout_index = (i + 1) % 3  # Cycle through 3 different layouts
            if layout_index == 0:
                content_slide_layout = prs.slide_layouts[1]  # Title and Content
            elif layout_index == 1:
                content_slide_layout = prs.slide_layouts[2]  # Section Header
            else:
                content_slide_layout = prs.slide_layouts[3]  # Two Content
            
            slide = prs.slides.add_slide(content_slide_layout)
            
            # Style the title
            if slide.shapes.title:
                title_shape = slide.shapes.title
                title_shape.text = slide_data["title"][0:255] if len(slide_data["title"]) > 255 else slide_data["title"]
                title_frame = title_shape.text_frame
                if title_frame.paragraphs:
                    title_para = title_frame.paragraphs[0]
                    title_para.font.size = Pt(32)
                    title_para.font.bold = True
                    title_para.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue
            
            # Handle content based on layout type
            content_text = slide_data["content"][0:1000] if len(slide_data["content"]) > 1000 else slide_data["content"]
            
            if layout_index == 0:  # Title and Content layout (layout 1)
                # Layout 1 typically has title (index 0) and content (index 1)
                if len(slide.placeholders) > 1:
                    content_placeholder = slide.placeholders[1]  # Usually the second placeholder
                    content_placeholder.text = content_text
                    
                    # Style the content text
                    content_frame = content_placeholder.text_frame
                    content_frame.word_wrap = True
                    
                    # Style paragraphs in content
                    for paragraph in content_frame.paragraphs:
                        paragraph.font.size = Pt(18)
                        paragraph.font.color.rgb = RGBColor(0, 0, 0)  # Black
                        paragraph.line_spacing = 1.3
            
            elif layout_index == 1:  # Section Header layout (layout 2)
                # Layout 2 typically has title (index 0) and subtitle (index 2)
                if len(slide.placeholders) > 2:
                    content_placeholder = slide.placeholders[2]  # Usually the subtitle placeholder
                    content_placeholder.text = content_text
                    
                    # Style the content text
                    content_frame = content_placeholder.text_frame
                    content_frame.word_wrap = True
                    
                    # Style paragraphs in content
                    for paragraph in content_frame.paragraphs:
                        paragraph.font.size = Pt(18)
                        paragraph.font.color.rgb = RGBColor(0, 0, 0)  # Black
                        paragraph.line_spacing = 1.3
            
            else:  # Two Content layout (layout 3)
                # Layout 3 typically has title (index 0), content1 (index 1), and content2 (index 2)
                if len(slide.placeholders) > 1:
                    content_placeholder = slide.placeholders[1]  # First content area
                    content_placeholder.text = content_text
                    
                    # Style the content text
                    content_frame = content_placeholder.text_frame
                    content_frame.word_wrap = True
                    
                    # Style paragraphs in content
                    for paragraph in content_frame.paragraphs:
                        paragraph.font.size = Pt(18)
                        paragraph.font.color.rgb = RGBColor(0, 0, 0)  # Black
                        paragraph.line_spacing = 1.3
                
                # Add secondary content to second placeholder if available
                if len(slide.placeholders) > 2:
                    secondary_placeholder = slide.placeholders[2]
                    secondary_placeholder.text = "Supporting details or examples related to the main content"
                    
                    secondary_frame = secondary_placeholder.text_frame
                    secondary_frame.word_wrap = True
                    
                    for paragraph in secondary_frame.paragraphs:
                        paragraph.font.size = Pt(16)
                        paragraph.font.color.rgb = RGBColor(50, 50, 50)
                        paragraph.line_spacing = 1.2
        
        # Add a thank you/conclusion slide
        thank_you_layout = prs.slide_layouts[0]
        thank_you_slide = prs.slides.add_slide(thank_you_layout)
        thank_title = thank_you_slide.shapes.title
        thank_subtitle = thank_you_slide.placeholders[1]
        
        if thank_title:
            thank_title.text = "Thank You"
            title_frame = thank_title.text_frame
            if title_frame.paragraphs:
                title_para = title_frame.paragraphs[0]
                title_para.font.size = Pt(44)
                title_para.font.bold = True
                title_para.font.color.rgb = RGBColor(0, 51, 102)
                title_para.alignment = PP_ALIGN.CENTER
        
        if thank_subtitle:
            thank_subtitle.text = "Questions & Discussion"
            subtitle_frame = thank_subtitle.text_frame
            if subtitle_frame.paragraphs:
                subtitle_para = subtitle_frame.paragraphs[0]
                subtitle_para.font.size = Pt(24)
                subtitle_para.font.color.rgb = RGBColor(102, 102, 102)
                subtitle_para.alignment = PP_ALIGN.CENTER
        
        # Generate filename
        safe_title = re.sub(r'[\\/*?:"<>|]', "", title)[:50]  # Remove invalid filename characters
        filename = os.path.join(self.output_dir, f"{safe_title}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx")
        
        # Save the presentation
        prs.save(filename)
        
        return filename
    
    def get_description(self) -> str:
        """Get skill description"""
        return "专业PPT制作工具，可以根据需求创建PowerPoint演示文稿，支持自动生成标题页和内容页"
